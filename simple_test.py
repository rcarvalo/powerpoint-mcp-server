#!/usr/bin/env python3
"""
Script de test simple et direct pour le PowerPoint MCP Server.
Ce script appelle directement les fonctions python-pptx pour vérifier que tout fonctionne.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

print("🚀 Test Simple - PowerPoint MCP Server")
print("="*60)

# Test 1: Créer une présentation
print("\n1️⃣  Création d'une nouvelle présentation...")
prs = Presentation()
print("✅ Présentation créée")

# Test 2: Ajouter des slides
print("\n2️⃣  Ajout de slides...")
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Test PowerPoint MCP Server"
subtitle.text = "Test Complet de Toutes les Fonctionnalités"
print("✅ Slide de titre ajouté")

# Test 3: Slide avec contenu
bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes = slide2.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Tests Réalisés'
tf = body_shape.text_frame
tf.text = 'Création de présentation'
p = tf.add_paragraph()
p.text = 'Ajout de slides'
p.level = 0
p = tf.add_paragraph()
p.text = 'Formatage de texte'
p.level = 1
print("✅ Slide avec bullets ajouté")

# Test 4: Ajouter une table
blank_slide_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(blank_slide_layout)
left = Inches(1.5)
top = Inches(2.0)
width = Inches(6.0)
height = Inches(2.5)

# Ajouter le titre
title_shape = slide3.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(0.5))
title_frame = title_shape.text_frame
title_frame.text = "Tableau de Données"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Ajouter la table
table = slide3.shapes.add_table(4, 3, left, top, width, height).table
table.cell(0, 0).text = 'Trimestre'
table.cell(0, 1).text = 'Revenus'
table.cell(0, 2).text = 'Profit'
table.cell(1, 0).text = 'Q1'
table.cell(1, 1).text = '100K€'
table.cell(1, 2).text = '20K€'
table.cell(2, 0).text = 'Q2'
table.cell(2, 1).text = '120K€'
table.cell(2, 2).text = '25K€'
table.cell(3, 0).text = 'Q3'
table.cell(3, 1).text = '150K€'
table.cell(3, 2).text = '35K€'
print("✅ Slide avec tableau ajouté")

# Test 5: Ajouter des formes
slide4 = prs.slides.add_slide(blank_slide_layout)

# Titre
title_shape = slide4.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(0.5))
title_frame = title_shape.text_frame
title_frame.text = "Formes et Éléments"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Rectangle
left = Inches(1.0)
top = Inches(1.5)
width = Inches(2.0)
height = Inches(1.5)
shape = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
shape.text = 'Rectangle'

# Oval
left = Inches(4.0)
shape = slide4.shapes.add_shape(
    MSO_SHAPE.OVAL, left, top, width, height
)
shape.text = 'Ovale'

# Triangle
left = Inches(7.0)
shape = slide4.shapes.add_shape(
    MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height
)
shape.text = 'Triangle'
print("✅ Slide avec formes ajouté")

# Test 6: Ajouter un graphique
slide5 = prs.slides.add_slide(blank_slide_layout)

# Titre
title_shape = slide5.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(0.5))
title_frame = title_shape.text_frame
title_frame.text = "Graphique de Performance"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Données du graphique
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('2023', (100, 120, 140, 160))
chart_data.add_series('2024', (130, 150, 170, 190))

# Ajouter le graphique
x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7.0), Inches(4.0)
chart = slide5.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.position = 2  # Right
print("✅ Slide avec graphique ajouté")

# Test 7: Ajouter une image (si disponible)
try:
    from PIL import Image
    import io

    # Créer une image simple pour le test
    img = Image.new('RGB', (300, 200), color=(73, 109, 137))
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)

    slide6 = prs.slides.add_slide(blank_slide_layout)

    # Titre
    title_shape = slide6.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(8.0), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "Image de Test"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Ajouter l'image
    left = Inches(2.5)
    top = Inches(2.0)
    pic = slide6.shapes.add_picture(img_bytes, left, top, width=Inches(5.0))
    print("✅ Slide avec image ajouté")
except Exception as e:
    print(f"⚠️  Image test skipped: {e}")

# Test 8: Propriétés du document
prs.core_properties.title = 'Test PowerPoint MCP Server'
prs.core_properties.author = 'PowerPoint MCP Test Suite'
prs.core_properties.subject = 'Tests Complets'
prs.core_properties.keywords = 'test, powerpoint, mcp, automation'
prs.core_properties.comments = 'Présentation générée automatiquement pour tester toutes les fonctionnalités'
print("\n✅ Propriétés du document définies")

# Test 9: Sauvegarder
output_file = 'test_simple_output.pptx'
prs.save(output_file)
print(f"\n✅ Présentation sauvegardée: {output_file}")

# Résumé
print("\n" + "="*60)
print("📊 RÉSUMÉ")
print("="*60)
print(f"✅ Nombre total de slides: {len(prs.slides)}")
print(f"✅ Fichier généré: {output_file}")
print("\n🎉 Tous les tests de base ont réussi!")
print("\nVous pouvez maintenant ouvrir le fichier test_simple_output.pptx")
print("="*60)
